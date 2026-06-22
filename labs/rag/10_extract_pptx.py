"""
10_extract_pptx.py

PPTX 문서에서 슬라이드별 텍스트를 추출한다.

입력:
- enterprise_docs/pptx/*.pptx

출력:
- extracted_text/pptx_extracted.jsonl

실행:
python 10_extract_pptx.py
"""

from __future__ import annotations

from pptx import Presentation

from step2_5_config import ENTERPRISE_DOCS_DIR, EXTRACTED_TEXT_DIR, ensure_directories
from step2_5_utils import clean_text, find_files, stable_id, write_jsonl, print_record_summary


def extract_shape_text(shape) -> list[str]:
    """PPT Shape에서 텍스트를 추출한다."""
    texts: list[str] = []

    # if hasattr(shape, "text"):
    #     text = clean_text(shape.text)
    #     if text:
    #         texts.append(text)

    # # 그룹 도형 내부 텍스트 처리
    # if hasattr(shape, "shapes"):
    #     for child in shape.shapes:
    #         texts.extend(extract_shape_text(child))

    # # 표 내부 텍스트 처리
    # if hasattr(shape, "table"):
    #     try:
    #         for row in shape.table.rows:
    #             row_values = []
    #             for cell in row.cells:
    #                 value = clean_text(cell.text)
    #                 if value:
    #                     row_values.append(value)
    #             if row_values:
    #                 texts.append(" | ".join(row_values))
    #     except Exception:
    #         pass


    # 일반 텍스트 도형 처리
    # has_text_frame이 True인 도형만 text를 읽는다.
    if getattr(shape, "has_text_frame", False):
        text = clean_text(shape.text)
        if text:
            texts.append(text)

    # 그룹 도형 내부 텍스트 처리
    # 그룹 도형인 경우에만 shapes 컬렉션을 순회한다.
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            texts.extend(extract_shape_text(child))

    # 표 내부 텍스트 처리
    # has_table이 True일 때만 shape.table에 접근한다.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_values = []

            for cell in row.cells:
                value = clean_text(cell.text)
                if value:
                    row_values.append(value)

            if row_values:
                texts.append(" | ".join(row_values))

    return texts


def extract_pptx_files() -> list[dict]:
    ensure_directories()

    pptx_dir = ENTERPRISE_DOCS_DIR / "pptx"
    pptx_files = find_files(pptx_dir, (".pptx",))

    records: list[dict] = []

    for pptx_path in pptx_files:
        try:
            presentation = Presentation(str(pptx_path))
        except Exception as exc:
            print(f"[PPTX 읽기 실패] {pptx_path.name}: {exc}")
            continue

        for slide_no, slide in enumerate(presentation.slides, start=1):
            slide_texts: list[str] = []

            for shape in slide.shapes:
                slide_texts.extend(extract_shape_text(shape))

            slide_content = clean_text("\n".join(slide_texts))

            if not slide_content:
                continue

            records.append(
                {
                    "id": stable_id(pptx_path.name, "pptx", slide_no),
                    "text": slide_content,
                    "metadata": {
                        "file_name": pptx_path.name,
                        "document_type": "pptx",
                        "slide_no": slide_no,
                        "source_path": str(pptx_path),
                        "extractor": "python-pptx",
                    },
                }
            )

    return records


def main() -> None:
    output_path = EXTRACTED_TEXT_DIR / "pptx_extracted.jsonl"
    records = extract_pptx_files()
    count = write_jsonl(output_path, records)
    print_record_summary("PPTX 텍스트 추출 완료", count, output_path)


if __name__ == "__main__":
    main()
